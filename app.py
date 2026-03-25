import io
import re

import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

# 1. 구글 시트 CSV 게시 링크
SHEET_CSV_URL = "https://docs.google.com/spreadsheets/d/e/2PACX-1vRhvBRSiCXgsdh9ipTVohLneYp55ROUHbSpOKUsQlVOJeuSB5bXbG7PDRWaB2VzrF57j2QPuy0G7Bcx/pub?output=csv"

# 2. 로컬에 저장된 PPT 양식 파일명
PPTX_TEMPLATE_PATH = "우리 팀의 오늘을 묻다_양식.pptx"
QUESTION_COUNT = 9
HIGHLIGHT_COLOR_HEX = "FFF59D"  # 연노랑 형광펜 색
SCORE_DISPLAY_ORDER = [4, 3, 2, 1]  # PPT 옵션 문단이 위에서부터 4,3,2,1 순서

@st.cache_data(ttl=60)
def load_data(url):
    return pd.read_csv(url)

def update_chart_data(chart, new_values):
    """기존 차트의 X축(Category)을 유지하면서 데이터(Value)만 교체합니다."""
    chart_data = CategoryChartData()
    
    categories = [c.label for c in chart.plots[0].categories]
    chart_data.categories = categories
    
    series_name = chart.series[0].name if chart.series else "Series 1"
    chart_data.add_series(series_name, new_values)
    
    chart.replace_data(chart_data)

def remove_unused_textboxes(slide, q1_score):
    """1번 문항 점수에 따라 해당하는 텍스트 상자 1개만 남기고 나머지를 삭제합니다."""
    if pd.isna(q1_score):
        return

    target_score = max(1, min(4, int(round(float(q1_score)))))
    emojis = {1: "🌧️", 2: "☁️", 3: "⛅", 4: "🌞"}
    target_emoji = emojis[target_score]

    shapes_to_delete = []
    seen_ids = set()

    for shape in slide.shapes:
        shape_id = shape.shape_id

        if shape.has_text_frame:
            text = shape.text.strip()
            if text in emojis.values():
                if text != target_emoji and shape_id not in seen_ids:
                    shapes_to_delete.append(shape)
                    seen_ids.add(shape_id)
                continue

        name = shape.name.lower().replace(" ", "")
        if "textbox1" in name or "텍스트상자1" in name:
            if target_score != 1 and shape_id not in seen_ids:
                shapes_to_delete.append(shape)
                seen_ids.add(shape_id)
        elif "textbox2" in name or "텍스트상자2" in name:
            if target_score != 2 and shape_id not in seen_ids:
                shapes_to_delete.append(shape)
                seen_ids.add(shape_id)
        elif "textbox3" in name or "텍스트상자3" in name:
            if target_score != 3 and shape_id not in seen_ids:
                shapes_to_delete.append(shape)
                seen_ids.add(shape_id)
        elif "textbox4" in name or "텍스트상자4" in name:
            if target_score != 4 and shape_id not in seen_ids:
                shapes_to_delete.append(shape)
                seen_ids.add(shape_id)

    for shape in shapes_to_delete:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
        except Exception:
            pass

def extract_question_number(value):
    match = re.match(r"\s*([1-9])\.", str(value))
    return int(match.group(1)) if match else None

def get_question_columns(df):
    question_columns = {}
    for col in df.columns:
        q_num = extract_question_number(col)
        if q_num and q_num not in question_columns:
            question_columns[q_num] = col
    return question_columns

def extract_rectangle_question_number(shape_name):
    match = re.search(r"(?:직사각형|rectangle)\s*(\d+)", str(shape_name), re.IGNORECASE)
    if not match:
        return None

    q_num = int(match.group(1))
    if 1 <= q_num <= QUESTION_COUNT:
        return q_num
    return None

def get_option_paragraph_indices(text_frame):
    """
    텍스트 프레임에서 선택지 문단 인덱스 4개를 PPT 기준으로 추출합니다.
    - 비어있지 않은 문단 중 마지막 4개를 선택지로 간주
    """
    paragraphs = list(text_frame.paragraphs)
    non_empty = [idx for idx, p in enumerate(paragraphs) if p.text and p.text.strip()]
    if len(non_empty) < 4:
        return []
    return non_empty[-4:]

def build_score_to_paragraph_index(text_frame):
    option_indices = get_option_paragraph_indices(text_frame)
    if len(option_indices) < 4:
        return {}
    return {
        score: paragraph_idx
        for score, paragraph_idx in zip(SCORE_DISPLAY_ORDER, option_indices)
    }

def set_run_highlight(run, color_hex=HIGHLIGHT_COLOR_HEX):
    r_pr = run._r.get_or_add_rPr()
    for child in list(r_pr):
        if child.tag == qn("a:highlight"):
            r_pr.remove(child)

    highlight = OxmlElement("a:highlight")
    color = OxmlElement("a:srgbClr")
    color.set("val", color_hex.upper())
    highlight.append(color)

    # PowerPoint는 <a:highlight> 위치(order)에 민감합니다.
    # 특히 uLnTx/uFillTx 앞에 들어가야 1~7번 텍스트 스타일에서도 하이라이트가 정상 렌더링됩니다.
    insert_before = {
        qn("a:uLnTx"),
        qn("a:uFillTx"),
        qn("a:latin"),
        qn("a:ea"),
        qn("a:cs"),
        qn("a:sym"),
        qn("a:hlinkClick"),
        qn("a:hlinkMouseOver"),
        qn("a:rtl"),
        qn("a:extLst"),
    }
    children = list(r_pr)
    insert_idx = next((idx for idx, child in enumerate(children) if child.tag in insert_before), None)
    if insert_idx is None:
        r_pr.append(highlight)
    else:
        r_pr.insert(insert_idx, highlight)

def highlight_paragraph(paragraph, color_hex=HIGHLIGHT_COLOR_HEX):
    if not paragraph.text or not paragraph.text.strip():
        return
    for run in paragraph.runs:
        if run.text and run.text.strip():
            set_run_highlight(run, color_hex)
            run.font.bold = True

def iter_all_shapes(shape_collection):
    for shape in shape_collection:
        yield shape
        if hasattr(shape, "shapes"):
            yield from iter_all_shapes(shape.shapes)

def apply_top_response_highlight(
    prs,
    top_scores_by_question,
):
    for slide in prs.slides:
        for shape in iter_all_shapes(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue

            q_num = extract_rectangle_question_number(shape.name)
            if not q_num:
                continue

            top_scores = top_scores_by_question.get(q_num, [])
            if not top_scores:
                continue

            score_to_paragraph = build_score_to_paragraph_index(shape.text_frame)
            if not score_to_paragraph:
                continue

            highlighted = set()
            for score in top_scores:
                paragraph_idx = score_to_paragraph.get(score)
                if paragraph_idx is None or paragraph_idx in highlighted:
                    continue
                highlight_paragraph(shape.text_frame.paragraphs[paragraph_idx])
                highlighted.add(paragraph_idx)

def main():
    st.title("설문 결과 PPT 자동 생성 툴")
    
    col1, col2 = st.columns([4, 1])
    with col1:
        st.write("원하는 차수(날짜)를 선택하고 PPT를 다운로드하세요.")
    with col2:
        # 데이터 새로고침 버튼
        if st.button("🔄 데이터 새로고침"):
            load_data.clear() # 캐시 초기화
            st.rerun() # 앱 재실행

    try:
        df = load_data(SHEET_CSV_URL)
    except Exception as e:
        st.error("구글 시트 데이터를 불러오는 데 실패했습니다. 링크를 확인해 주세요.")
        return

    date_column = '날짜' 
    if date_column not in df.columns:
        st.error(f"데이터에 '{date_column}' 컬럼이 존재하지 않습니다.")
        return
        
    dates = df[date_column].dropna().unique()
    selected_date = st.selectbox("출력할 날짜(차수)를 선택하세요:", sorted(dates))
    
    if selected_date:
        filtered_df = df[df[date_column] == selected_date]

        question_columns = get_question_columns(df)
        missing_questions = [q for q in range(1, QUESTION_COUNT + 1) if q not in question_columns]
        if missing_questions:
            st.error("1번부터 9번까지의 문항 컬럼을 모두 찾을 수 없습니다. 시트의 열 제목을 확인해 주세요.")
            return

        ordered_q_cols = [question_columns[q] for q in range(1, QUESTION_COUNT + 1)]
        numeric_answers = filtered_df[ordered_q_cols].apply(pd.to_numeric, errors="coerce")
        averages = numeric_answers.mean().round(1).fillna(0).tolist()

        top_scores_by_question = {}
        top_frequency_by_question = {}
        for idx, col_name in enumerate(ordered_q_cols, start=1):
            counts = numeric_answers[col_name].dropna().astype(int).value_counts()
            if counts.empty:
                top_scores_by_question[idx] = []
                top_frequency_by_question[idx] = 0
                continue

            max_freq = int(counts.max())
            top_scores = sorted(int(score) for score in counts[counts == max_freq].index.tolist())
            top_scores_by_question[idx] = top_scores
            top_frequency_by_question[idx] = max_freq

        st.write(f"**{selected_date}** 평균 데이터 미리보기:")
        st.dataframe(pd.DataFrame([averages], columns=[f"Q{i}" for i in range(1, QUESTION_COUNT + 1)]))

        mode_rows = []
        for q in range(1, QUESTION_COUNT + 1):
            top_scores = top_scores_by_question[q]
            mode_rows.append(
                {
                    "문항": f"Q{q}",
                    "최빈 응답": ", ".join(map(str, top_scores)) if top_scores else "-",
                    "빈도수": top_frequency_by_question[q],
                }
            )
        st.write(f"**{selected_date}** 최빈 응답(빈도수):")
        st.dataframe(pd.DataFrame(mode_rows))

        if st.button("PPT 생성 및 다운로드 준비"):
            try:
                prs = Presentation(PPTX_TEMPLATE_PATH)
            except Exception as e:
                st.error(f"PPT 양식 파일을 찾을 수 없습니다. '{PPTX_TEMPLATE_PATH}' 파일이 같은 폴더에 있는지 확인해 주세요.")
                return
            
            # --- 1페이지 (Index 0) ---
            if len(prs.slides) >= 1:
                remove_unused_textboxes(prs.slides[0], averages[0])
                for shape in prs.slides[0].shapes:
                    if shape.has_chart:
                        update_chart_data(shape.chart, averages[0:3])
                        break
                    
            # --- 2페이지 (Index 1) ---
            if len(prs.slides) >= 2:
                for shape in prs.slides[1].shapes:
                    if shape.has_chart:
                        update_chart_data(shape.chart, averages[3:6])
                        break
                    
            # --- 3페이지 (Index 2) ---
            if len(prs.slides) >= 3:
                for shape in prs.slides[2].shapes:
                    if shape.has_chart:
                        update_chart_data(shape.chart, averages[6:9])
                        break
                    
            # --- 4페이지 (Index 3) ---
            if len(prs.slides) >= 4:
                remove_unused_textboxes(prs.slides[3], averages[0])
                chart_idx = 0
                for shape in prs.slides[3].shapes:
                    if shape.has_chart:
                        if chart_idx == 0:
                            update_chart_data(shape.chart, averages[0:3])
                        elif chart_idx == 1:
                            update_chart_data(shape.chart, averages[3:6])
                        elif chart_idx == 2:
                            update_chart_data(shape.chart, averages[6:9])
                        chart_idx += 1
            
            # 직사각형 1~9에서 최빈 응답 선택지에 형광펜 표시
            apply_top_response_highlight(
                prs,
                top_scores_by_question,
            )

            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            
            st.success("PPT 생성이 완료되었습니다.")
            st.download_button(
                label="📥 PPT 다운로드",
                data=output,
                file_name=f"팀_설문결과_리포트_{selected_date}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

if __name__ == "__main__":
    main()
